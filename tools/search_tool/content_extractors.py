"""
Extracción de contenido de archivos: DOCX, PDF, XLSX, PPTX, TXT.
Separado de processor.py por SRP (máxima R0: clases <300 líneas).
"""
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# Libraries para extraer contenido
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def extract_docx_content(file_path: str) -> str:
    """Extrae texto de un archivo DOCX."""
    if not DOCX_AVAILABLE:
        return ""
    
    try:
        doc = DocxDocument(file_path)
        text = ' '.join([p.text for p in doc.paragraphs if p.text])
        return text
    except Exception as e:
        logger.warning(f"Error extracting DOCX: {e}")
        return ""


def extract_pdf_content(file_path: str) -> str:
    """Extrae texto de un archivo PDF."""
    if not PDF_AVAILABLE:
        return ""
    
    try:
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + " "
        return text
    except Exception as e:
        logger.warning(f"Error extracting PDF: {e}")
        return ""


def extract_xlsx_content(file_path: str) -> str:
    """Extrae texto de un archivo XLSX."""
    if not XLSX_AVAILABLE:
        return ""
    
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value:
                        text += str(cell.value) + " "
        return text
    except Exception as e:
        logger.warning(f"Error extracting XLSX: {e}")
        return ""


def extract_pptx_content(file_path: str) -> str:
    """Extrae texto de un archivo PPTX."""
    if not PPTX_AVAILABLE:
        return ""
    
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text
    except Exception as e:
        logger.warning(f"Error extracting PPTX: {e}")
        return ""


def get_file_content(file_path: str) -> str:
    """Extrae contenido según tipo de archivo."""
    ext = Path(file_path).suffix.lower()
    
    if ext == '.docx' or ext == '.doc':
        return extract_docx_content(file_path)
    elif ext == '.pdf':
        return extract_pdf_content(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_xlsx_content(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_pptx_content(file_path)
    elif ext == '.txt':
        try:
            with open(file_path, 'r', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.debug(f"Error reading txt: {e}")
            return ""
    
    return ""


def search_content(files: list, pattern: str, case_sensitive: bool = False) -> dict:
    """Busca contenido en archivos."""
    results = {}
    
    for f in files:
        content = get_file_content(f)
        if not content:
            continue
        
        if case_sensitive:
            if pattern in content:
                count = content.count(pattern)
                results[f] = {'matches': count, 'content': content[:500]}
        else:
            lower_content = content.lower()
            lower_pattern = pattern.lower()
            if lower_pattern in lower_content:
                count = lower_content.count(lower_pattern)
                results[f] = {'matches': count, 'content': content[:500]}
    
    return results