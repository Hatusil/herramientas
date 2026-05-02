"""
Processor: Funciones de búsqueda avanzada de archivos.
"""
import logging

logger = logging.getLogger(__name__)
import os
import re
import csv
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional

# Flag para cancelar búsqueda
SEARCH_CANCELLED = False

def cancel_search():
    """Cancela la búsqueda en curso."""
    global SEARCH_CANCELLED
    SEARCH_CANCELLED = True

def reset_search():
    """Resetea el flag de cancelación."""
    global SEARCH_CANCELLED
    SEARCH_CANCELLED = False

# Libraries para extraer contenido
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def search_by_name(files: List[str], pattern: str, mode: str = 'contains', 
                  case_sensitive: bool = False) -> List[str]:
    """Busca archivos por nombre."""
    results = []
    
    if not pattern:
        return files
    
    for f in files:
        name = os.path.basename(f)
        
        if mode == 'exact':
            match = (name == pattern) if case_sensitive else (name.lower() == pattern.lower())
        elif mode == 'contains':
            match = (pattern in name) if case_sensitive else (pattern.lower() in name.lower())
        elif mode == 'regex':
            try:
                flags = 0 if case_sensitive else re.IGNORECASE
                match = bool(re.search(pattern, name, flags))
            except re.error:
                match = False
        else:
            match = pattern.lower() in name.lower()
        
        if match:
            results.append(f)
    
    return results


def search_by_date(files: List[str], date_from: Optional[str] = None,
                   date_to: Optional[str] = None) -> List[str]:
    """Filtra archivos por fecha de modificación."""
    results = []
    
    # Parse dates
    from_date = None
    to_date = None
    
    if date_from:
        try:
            from_date = datetime.strptime(date_from, '%d/%m/%Y')
        except ValueError:
            try:
                from_date = datetime.strptime(date_from, '%Y-%m-%d')
            except ValueError as e:
                logger.warning(f"Invalid date from format: {e}")
    
    if date_to:
        try:
            to_date = datetime.strptime(date_to, '%d/%m/%Y')
        except ValueError:
            try:
                to_date = datetime.strptime(date_to, '%Y-%m-%d')
            except ValueError as e:
                logger.warning(f"Invalid date from format: {e}")
    
    for f in files:
        if not os.path.exists(f):
            continue
        
        mtime = datetime.fromtimestamp(os.path.getmtime(f))
        
        if from_date and mtime < from_date:
            continue
        if to_date and mtime > to_date:
            continue
        
        results.append(f)
    
    return results


def filter_by_size(files: List[str], min_size: Optional[int] = None,
                  max_size: Optional[int] = None) -> List[str]:
    """Filtra archivos por tamaño en bytes."""
    results = []
    
    for f in files:
        if not os.path.exists(f):
            continue
        
        size = os.path.getsize(f)
        
        if min_size and size < min_size:
            continue
        if max_size and size > max_size:
            continue
        
        results.append(f)
    
    return results


def filter_by_extension(files: List[str], extensions: List[str]) -> List[str]:
    """Filtra por extensión."""
    if not extensions:
        return files
    
    extensions = [e.lower().replace('.', '') for e in extensions]
    
    results = []
    for f in files:
        ext = Path(f).suffix.lower().replace('.', '')
        if ext in extensions:
            results.append(f)
    
    return results


def extract_docx_content(file_path: str) -> str:
    """Extrae texto de un archivo DOCX."""
    if not DOCX_AVAILABLE:
        return ""
    
    try:
        doc = DocxDocument(file_path)
        text = ' '.join([p.text for p in doc.paragraphs if p.text])
        return text
    except Exception as e:
        logger.warning(f"Error extracting: {e}")
        return ""


def extract_pdf_content(file_path: str) -> str:
    """Extrae texto de un archivo PDF."""
    if not PDF_AVAILABLE:
        return ""
    
    try:
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + " "
        return text
    except Exception as e:
        logger.warning(f"Error extracting: {e}")
        return ""


def extract_xlsx_content(file_path: str) -> str:
    """Extrae texto de un archivo XLSX."""
    if not XLSX_AVAILABLE:
        return ""
    
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value:
                        text += str(cell.value) + " "
        return text
    except Exception as e:
        logger.warning(f"Error extracting: {e}")
        return ""


def extract_pptx_content(file_path: str) -> str:
    """Extrae texto de un archivo PPTX."""
    if not PPTX_AVAILABLE:
        return ""
    
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text
    except Exception as e:
        logger.warning(f"Error extracting: {e}")
        return ""


def get_file_content(file_path: str) -> str:
    """Extrae contenido según tipo de archivo."""
    ext = Path(file_path).suffix.lower()
    
    if ext == '.docx' or ext == '.doc':
        return extract_docx_content(file_path)
    elif ext == '.pdf':
        return extract_pdf_content(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_xlsx_content(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_pptx_content(file_path)
    elif ext == '.txt':
        try:
            with open(file_path, 'r', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.debug(f"Error reading txt: {e}")
            return ""
    
    return ""


def search_content(files: List[str], pattern: str, case_sensitive: bool = False) -> Dict[str, Any]:
    """Busca contenido en archivos."""
    results = {}
    
    for f in files:
        content = get_file_content(f)
        if not content:
            continue
        
        if case_sensitive:
            if pattern in content:
                # Contar ocurrencias
                count = content.count(pattern)
                results[f] = {'matches': count, 'content': content[:500]}
        else:
            lower_content = content.lower()
            lower_pattern = pattern.lower()
            if lower_pattern in lower_content:
                count = lower_content.count(lower_pattern)
                results[f] = {'matches': count, 'content': content[:500]}
    
    return results


def search_all(folder: str, options: Dict[str, Any]) -> Dict[str, Any]:
    """
    Función principal de búsqueda.
    
    options: {
        'name_pattern': str,
        'name_mode': 'exact'|'contains'|'regex',
        'case_sensitive': bool,
        'date_from': str|null,
        'date_to': str|null,
        'extensions': list,
        'min_size': int|null,
        'max_size': int|null,
        'search_content': bool,
        'content_pattern': str,
    }
    """
    global SEARCH_CANCELLED
    
    if not os.path.exists(folder):
        return {'success': False, 'error': 'Carpeta no encontrada'}
    
    # Recolectar archivos recursivamente - con verificación frecuente
    all_files = []
    file_count = 0
    for root, dirs, files in os.walk(folder):
        # Verificar cancelación frecuentemente (cada 50 archivos)
        file_count += len(files)
        if file_count >= 50:
            if SEARCH_CANCELLED:
                return {'success': True, 'cancelled': True, 'results': [], 'count': 0}
            file_count = 0
        
        for name in files:
            all_files.append(os.path.join(root, name))
    
    # Verificar cancelación después de recolectar
    if SEARCH_CANCELLED:
        return {'success': True, 'cancelled': True, 'results': [], 'count': 0}
    
    # Aplicar filtros
    results = all_files
    
    # Por nombre
    if options.get('name_pattern'):
        results = search_by_name(
            results,
            options['name_pattern'],
            options.get('name_mode', 'contains'),
            options.get('case_sensitive', False)
        )
        
        # Verificar cancelación después de cada filtro
        if SEARCH_CANCELLED:
            return {'success': True, 'cancelled': True, 'results': [], 'count': 0}
    
    # Por fecha
    if options.get('date_from') or options.get('date_to'):
        results = search_by_date(results, options.get('date_from'), options.get('date_to'))
    
    # Por tamaño
    if options.get('min_size') or options.get('max_size'):
        results = filter_by_size(results, options.get('min_size'), options.get('max_size'))
    
    # Por extensión
    if options.get('extensions'):
        results = filter_by_extension(results, options['extensions'])
    
    # Buscar contenido (puede tardar mucho - verificar frecuentemente)
    content_results = {}
    if options.get('search_content') and options.get('content_pattern'):
        pattern = options['content_pattern']
        case_sensitive = options.get('case_sensitive', False)
        
        for i, f in enumerate(results):
            if i % 10 == 0:  # Cada 10 archivos verificar
                if SEARCH_CANCELLED:
                    return {'success': True, 'cancelled': True, 'results': [], 'count': 0}
            
            content = get_file_content(f)
            if not content:
                continue
            
            if case_sensitive:
                if pattern in content:
                    count = content.count(pattern)
                    content_results[f] = {'matches': count, 'content': content[:500]}
            else:
                lower_content = content.lower()
                lower_pattern = pattern.lower()
                if lower_pattern in lower_content:
                    count = lower_content.count(lower_pattern)
                    content_results[f] = {'matches': count, 'content': content[:500]}
    
    # Verificar cancelación final
    if SEARCH_CANCELLED:
        return {'success': True, 'cancelled': True, 'results': [], 'count': 0}
    
    # Si busca contenido Y hay matches, filtrar solo archivos con contenido
    if options.get('search_content') and options.get('content_pattern') and content_results:
        results = list(content_results.keys())
    
    # Preparar resultados finales
    final_results = []
    for f in results:
        is_file = os.path.isfile(f)
        stat = os.stat(f) if is_file else None
        
        result = {
            'path': f,
            'name': os.path.basename(f),
            'size': stat.st_size if stat else 0,
            'modified': datetime.fromtimestamp(stat.st_mtime).strftime('%d/%m/%Y %H:%M') if stat else '',
            'matches': content_results.get(f, {}).get('matches', 0) if f in content_results else 0
        }
        final_results.append(result)
    
    return {
        'success': True,
        'results': final_results,
        'count': len(final_results),
        'content_matches': content_results
    }


def export_to_csv(results: List[Dict], output_path: str) -> bool:
    """Exporta resultados a CSV."""
    try:
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.DictWriter(f, fieldnames=['path', 'name', 'size', 'modified', 'matches'])
            writer.writeheader()
            writer.writerows(results)
        return True
    except Exception as e:
        logger.error(f"Export CSV error: {e}")
        return False


def export_to_txt(results: List[Dict], output_path: str) -> bool:
    """Exporta resultados a TXT."""
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            for r in results:
                f.write(f"{r['path']}\n")
        return True
    except Exception as e:
        logger.error(f"Export TXT error: {e}")
        return False